# coding: utf-8
# Copyright (c) Huawei Technologies Co., Ltd. 2026. All rights reserved.

import base64
import html
import io
import json
import os
import pathlib
import re
import shlex
import sys
import asyncio
import shutil
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from typing import Dict, Any, AsyncIterator, List, Optional, Tuple

from openjiuwen.core.common.exception.codes import StatusCode
from openjiuwen.core.common.logging import logger
from openjiuwen.core.foundation.tool.base import Tool
from openjiuwen.core.sys_operation import OperationMode, SysOperation
from openjiuwen.core.sys_operation.cwd import get_cwd
from openjiuwen.harness.prompts.sections.tools import build_tool_card
from openjiuwen.harness.tools.base_tool import ToolOutput

# Device files that would block indefinitely or produce infinite output.
_BLOCKED_DEVICE_PATHS: frozenset = frozenset({
    "/dev/zero", "/dev/random", "/dev/urandom", "/dev/full",
    "/dev/stdin", "/dev/tty", "/dev/console",
    "/dev/stdout", "/dev/stderr",
    "/dev/fd/0", "/dev/fd/1", "/dev/fd/2",
})

# Binary file extensions that are not directly readable as text.
# PDF and image extensions are handled separately and are NOT in this set.
_BINARY_EXTENSIONS: frozenset = frozenset({
    ".exe", ".dll", ".so", ".dylib", ".bin", ".obj", ".o", ".a", ".lib",
    ".zip", ".tar", ".gz", ".bz2", ".xz", ".7z", ".rar",
    ".pyc", ".pyo", ".class", ".wasm",
    ".db", ".sqlite", ".sqlite3",
})

# Office document extensions (ZIP-based binary containers, need dedicated parsers).
_OFFICE_DOC_EXTENSIONS: frozenset = frozenset({
    ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
})


@dataclass
class _FileReadState:
    """State recorded when a file is read; used by EditFileTool for pre-read validation."""
    mtime_ns: int
    size_bytes: int
    is_partial: bool  # True if only a subset of lines was read (offset > 0 or explicit limit)
    content: Optional[str] = None


@dataclass
class _RawTextState:
    content: str
    line_count: int


# Module-level read state registry.
# ReadFileTool populates it on successful text reads.
# EditFileTool consumes it to enforce "must read before edit" and detect external modifications.
_FILE_READ_REGISTRY: Dict[str, _FileReadState] = {}

_HISTORY_LOCK = asyncio.Lock()
_PDF_READ_SEMAPHORE = asyncio.Semaphore(1)
MAX_HISTORY_PER_FILE: int = 100


async def _append_op_history(history_path: str, file_path: str, action: str,
                              old_content: Optional[str], new_content: str) -> None:
    """Append a write/edit operation to the per-workspace history JSON file.

    Async-safe: concurrent coroutines share _HISTORY_LOCK so reads and writes
    to the JSON file are serialised without blocking the event loop.
    The JSON file is the single source of truth; there is no in-memory cache.
    """
    entry = {
        "action": action,
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "old_content": old_content,
        "new_content": new_content,
    }
    try:
        async with _HISTORY_LOCK:
            history: Dict[str, list] = {}
            if os.path.exists(history_path):
                with open(history_path, "r", encoding="utf-8") as f:
                    history = json.load(f)
            entries = history.setdefault(file_path, [])
            entries.append(entry)
            if len(entries) > MAX_HISTORY_PER_FILE:
                history[file_path] = entries[-MAX_HISTORY_PER_FILE:]
            os.makedirs(os.path.dirname(history_path), exist_ok=True)
            tmp = history_path + ".tmp"
            with open(tmp, "w", encoding="utf-8") as f:
                json.dump(history, f, ensure_ascii=False, indent=2)
            os.replace(tmp, history_path)
    except Exception as exc:  # noqa: BLE001
        logger.warning("[_append_op_history] Failed to persist file op history to %s: %s", history_path, exc)


class MaxFileReadTokenExceededError(Exception):
    """Raised when file content exceeds the maximum allowed token count."""

    def __init__(self, token_count: int, max_tokens: int) -> None:
        super().__init__(
            f"File content ({token_count} tokens) exceeds maximum allowed tokens ({max_tokens}). "
            "Use offset and limit parameters to read specific portions of the file, "
            "or search for specific content instead of reading the whole file."
        )
        self.token_count = token_count
        self.max_tokens = max_tokens


class PDFReadError(RuntimeError):
    """Raised for structured PDF read errors that the model can recover from."""


@dataclass
class _ReadSnapshot:
    mtime_ns: int
    line_count: int


def _resolve_tool_file_path(operation: SysOperation, file_path: str) -> str:
    """Resolve relative tool paths against the configured sys_operation work_dir.

    Keeps UNC paths unchanged. Relative paths are only accepted when the operation
    exposes a work_dir; otherwise the caller must still provide an absolute path.
    """
    expanded = os.path.expanduser(file_path)
    if expanded.startswith("\\\\") or expanded.startswith("//") or os.path.isabs(expanded):
        return expanded

    work_dir = get_cwd()
    return str((pathlib.Path(work_dir).expanduser().resolve() / expanded).resolve())


def _is_unc_path(path_value: str) -> bool:
    return path_value.startswith("\\\\") or path_value.startswith("//")


class ReadFileTool(Tool):
    MAX_LINES_TO_READ: int = 2000
    MAX_SIZE_BYTES: int = 256 * 1024   # 256 KB — for text files (offset/limit can split)
    MAX_TOKENS: int = 25_000
    MAX_PDF_SIZE_BYTES_WITHOUT_PAGES: int = 10 * 1024 * 1024
    MAX_PDF_ABSOLUTE_SIZE_BYTES: int = 200 * 1024 * 1024
    MAX_OFFICE_DOC_SIZE_BYTES: int = 10 * 1024 * 1024  # 10 MB — ZIP containers with media
    PDF_MAX_PAGES_PER_READ: int = 10
    PDF_AT_MENTION_INLINE_THRESHOLD: int = 10
    FILE_UNCHANGED_STUB: str = "File unchanged since last read. Reuse the previously returned content."
    _IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff"}

    def __init__(
        self,
        operation: SysOperation,
        language: str = "cn",
        agent_id: Optional[str] = None,
        enable_image_multimodal: bool = True,
    ):
        super().__init__(build_tool_card("read_file", "ReadFileTool", language, agent_id=agent_id))
        self.operation = operation
        self._snapshots: Dict[Tuple[str, int, int, str], _ReadSnapshot] = {}
        self.enable_image_multimodal = enable_image_multimodal

    # ------------------------------------------------------------------
    # File-type predicates
    # ------------------------------------------------------------------

    @staticmethod
    def _is_pdf(file_path: str) -> bool:
        return file_path.lower().endswith(".pdf")

    @classmethod
    def _is_image(cls, file_path: str) -> bool:
        _, ext = os.path.splitext(file_path.lower())
        return ext in cls._IMAGE_EXTENSIONS

    @staticmethod
    def _is_pdf_supported(model_name: str) -> bool:
        """Returns False for claude-3-haiku variants, which cannot process inline PDFs."""
        return "claude-3-haiku" not in model_name.lower()

    @staticmethod
    def _is_blocked_device(file_path: str) -> bool:
        if file_path in _BLOCKED_DEVICE_PATHS:
            return True
        # Linux /proc/<pid>/fd/0-2 aliases for stdio
        if file_path.startswith("/proc/") and file_path.endswith(("/fd/0", "/fd/1", "/fd/2")):
            return True
        return False

    @staticmethod
    def _is_binary(file_path: str) -> bool:
        _, ext = os.path.splitext(file_path.lower())
        return ext in _BINARY_EXTENSIONS

    @classmethod
    def _is_office_doc(cls, file_path: str) -> bool:
        _, ext = os.path.splitext(file_path.lower())
        return ext in _OFFICE_DOC_EXTENSIONS

    def _is_exempt_from_binary_check(self, file_path: str) -> bool:
        """Whether a file has a dedicated reader and is exempt from the binary-text guard."""
        return (
            self._is_pdf(file_path)
            or self._is_image(file_path)
            or self._is_office_doc(file_path)
            or self._is_plain_text_candidate(file_path)
        )

    # ------------------------------------------------------------------
    # PDF page range parsing
    # ------------------------------------------------------------------

    @staticmethod
    def _validate_pdf_page_range_format(pages: str) -> bool:
        """Return True when *pages* is syntactically valid (ignores document page count)."""
        raw = str(pages).strip()
        if not raw:
            return False
        try:
            if "-" not in raw:
                return int(raw) >= 1
            start_text, end_text = raw.split("-", 1)
            if start_text and int(start_text) < 1:
                return False
            if end_text and int(end_text) < 1:
                return False
            return True
        except (ValueError, TypeError):
            return False

    @staticmethod
    def _parse_pdf_page_range_with_reason(
            pages: Optional[str],
            total_pages: int,
    ) -> Tuple[Optional[Tuple[int, int]], Optional[str]]:
        """Parse *pages* into an inclusive (start, end) pair.

        Returns (range, None) on success, or (None, reason) where reason is
        ``invalid_format`` or ``out_of_bounds``.
        """
        if not pages:
            if total_pages <= 0:
                return None, "invalid_format"
            return (1, total_pages), None

        if not ReadFileTool._validate_pdf_page_range_format(pages):
            return None, "invalid_format"

        raw = str(pages).strip()
        try:
            if "-" not in raw:
                page = int(raw)
                if page > total_pages:
                    return None, "out_of_bounds"
                return (page, page), None
            start_text, end_text = raw.split("-", 1)
            start = int(start_text) if start_text else 1
            end = int(end_text) if end_text else total_pages
            if start > total_pages:
                return None, "out_of_bounds"
            end = min(total_pages, end)
            if start > end:
                return None, "out_of_bounds"
            return (start, end), None
        except (ValueError, TypeError):
            return None, "invalid_format"

    @staticmethod
    def _parse_pdf_page_range(pages: Optional[str], total_pages: int) -> Optional[Tuple[int, int]]:
        """Parse a PDF page range string into an inclusive (start, end) pair.

        Accepts formats: "3", "1-5", "10-", "-5".
        Returns None when the range string is malformed or yields an empty range.
        """
        result, _reason = ReadFileTool._parse_pdf_page_range_with_reason(pages, total_pages)
        return result

    # ------------------------------------------------------------------
    # Capability flags
    # ------------------------------------------------------------------

    @staticmethod
    def is_read_only() -> bool:
        return True

    @staticmethod
    def is_concurrency_safe() -> bool:
        return True

    @staticmethod
    def check_permissions() -> str:
        return "allow"

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _resolve_model_name(kwargs: Dict[str, Any]) -> str:
        return str(
            kwargs.get("model")
            or kwargs.get("model_name")
            or kwargs.get("llm_model")
            or ""
        )

    @staticmethod
    def _cat_n(text: str) -> str:
        """Format text with cat-style 1-indexed line numbers."""
        if text == "":
            return ""
        lines = text.splitlines()
        return "\n".join(f"{idx:>6}\t{line}" for idx, line in enumerate(lines, 1))

    def _estimate_tokens(self, text: str) -> int:
        """Rough token estimate: ~4 UTF-8 characters per token."""
        return max(1, len(text) // 4)

    @staticmethod
    def _raw_line_count(text: str) -> int:
        return len(text.splitlines()) if text else 0

    @staticmethod
    def _format_mb(size_bytes: int) -> int:
        return max(1, size_bytes // (1024 * 1024))

    @staticmethod
    def _suggest_page_batches(max_pages: int = 30, batch_size: int = 10) -> List[str]:
        ranges = []
        start = 1
        while start <= max_pages:
            end = min(max_pages, start + batch_size - 1)
            ranges.append(f"{start}-{end}")
            start = end + 1
        return ranges

    def _build_pdf_too_large_no_pages_error(self, file_path: str, size_bytes: int) -> str:
        batches = self._suggest_page_batches(max_pages=30, batch_size=self.PDF_MAX_PAGES_PER_READ)
        examples = "\n".join(
            f'  {idx}. read_file(file_path="{file_path}", pages="{page_range}")'
            for idx, page_range in enumerate(batches[:2], 1)
        )
        return (
            "[PDF_READ_ERROR] CODE=PDF_TOO_LARGE_NO_PAGES\n"
            f"File: {file_path} ({self._format_mb(size_bytes)} MB)\n"
            f"Reason: Files over {self._format_mb(self.MAX_PDF_SIZE_BYTES_WITHOUT_PAGES)} MB "
            "cannot be read without the pages parameter.\n\n"
            "ACTION REQUIRED — retry with paginated reads:\n"
            f"{examples}\n"
            f'  3. If more content is needed, continue with pages="{batches[2]}", then "31-40", ...\n'
            "  4. Stop when read_file returns PDF_PAGE_RANGE_OUT_OF_BOUNDS, then summarize all collected batches.\n"
            "Do NOT call read_file without pages on this file again."
        )

    def _build_pdf_too_many_pages_error(self, file_path: str, total_pages: int) -> str:
        first = f"1-{min(self.PDF_MAX_PAGES_PER_READ, total_pages)}"
        return (
            "[PDF_READ_ERROR] CODE=PDF_TOO_MANY_PAGES_NO_PAGES\n"
            f"File: {file_path} ({total_pages} pages)\n"
            f"Reason: PDFs with more than {self.PDF_AT_MENTION_INLINE_THRESHOLD} pages require the pages parameter.\n\n"
            "ACTION REQUIRED — retry with paginated reads:\n"
            f'  read_file(file_path="{file_path}", pages="{first}")'
        )

    def _build_pdf_page_range_too_wide_error(self, file_path: str, pages: str, page_count: int) -> str:
        return (
            "[PDF_READ_ERROR] CODE=PDF_PAGE_RANGE_TOO_WIDE\n"
            f"File: {file_path}\n"
            f'Requested: pages="{pages}" ({page_count} pages)\n'
            f"Limit: {self.PDF_MAX_PAGES_PER_READ} pages per request.\n\n"
            "ACTION REQUIRED — split into smaller ranges, for example:\n"
            f'  read_file(file_path="{file_path}", pages="1-{self.PDF_MAX_PAGES_PER_READ}")'
        )

    def _build_pdf_page_range_out_of_bounds_error(
            self, file_path: str, pages: Optional[str], total_pages: int
    ) -> str:
        return (
            "[PDF_READ_ERROR] CODE=PDF_PAGE_RANGE_OUT_OF_BOUNDS\n"
            f"File: {file_path} ({total_pages} pages)\n"
            f"Requested: pages={pages!r}\n"
            "Reason: The requested page range is outside the PDF page count.\n\n"
            "ACTION REQUIRED — stop reading additional page batches and summarize the content already collected."
        )

    def _build_pdf_absolute_size_error(self, file_path: str, size_bytes: int, pages: Optional[str]) -> str:
        return (
            "[PDF_READ_ERROR] CODE=PDF_ABSOLUTE_SIZE_EXCEEDED\n"
            f"File: {file_path} ({self._format_mb(size_bytes)} MB)\n"
            f"Requested pages: {pages!r}\n"
            f"Limit: {self._format_mb(self.MAX_PDF_ABSOLUTE_SIZE_BYTES)} MB absolute maximum.\n\n"
            "ACTION REQUIRED — do NOT retry read_file on this file. Ask the user to split the PDF "
            "or use a dedicated PDF/OCR workflow."
        )

    def _build_pdf_sandbox_large_file_error(self, file_path: str, pages: Optional[str]) -> str:
        return (
            "[PDF_READ_ERROR] CODE=PDF_SANDBOX_LARGE_FILE_UNSUPPORTED\n"
            f"File: {file_path}\n"
            f"Requested pages: {pages!r}\n"
            f"Reason: SANDBOX mode only supports PDF files up to "
            f"{self._format_mb(self.MAX_PDF_SIZE_BYTES_WITHOUT_PAGES)} MB. "
            "Paginated reads are allowed within that size limit.\n\n"
            "ACTION REQUIRED — use LOCAL mode or a dedicated remote PDF parsing workflow "
            "for larger files."
        )

    def _build_pdf_invalid_page_range_error(self, pages: str) -> str:
        return (
            "[PDF_READ_ERROR] CODE=PDF_INVALID_PAGE_RANGE\n"
            f"Invalid PDF page range format: '{pages}'. "
            "Use formats like '3' or '1-10'. Pages are 1-indexed."
        )

    def _suggest_reduced_pdf_page_range(self, pages: Optional[str]) -> Tuple[str, bool]:
        """Suggest a narrower pages range after token overflow.

        Returns (suggested_pages, single_page_unreducible).
        """
        fallback = f"1-{max(1, self.PDF_MAX_PAGES_PER_READ // 2)}"
        if not pages or not str(pages).strip():
            return fallback, False

        parsed = self._parse_pdf_page_range(str(pages).strip(), sys.maxsize)
        if parsed is None:
            return fallback, False

        start, end = parsed
        page_count = end - start + 1
        if page_count <= 1:
            return str(start), True

        reduced_count = max(1, (page_count + 1) // 2)
        new_end = start + reduced_count - 1
        if start == new_end:
            return str(start), False
        return f"{start}-{new_end}", False

    def _build_pdf_output_token_error(
            self, file_path: str, pages: Optional[str], token_count: int
    ) -> str:
        suggested_pages, single_page_unreducible = self._suggest_reduced_pdf_page_range(pages)
        if single_page_unreducible:
            action = (
                "ACTION REQUIRED — this single page exceeds the token limit and cannot be "
                "split further via pages. Ask the user to split the PDF or use a dedicated OCR workflow."
            )
        else:
            action = (
                "ACTION REQUIRED — retry with a smaller page range:\n"
                f'  read_file(file_path="{file_path}", pages="{suggested_pages}")'
            )
        return (
            "[PDF_READ_ERROR] CODE=PDF_OUTPUT_TOKEN_EXCEEDED\n"
            f"File: {file_path}, pages={pages!r}\n"
            f"Reason: Extracted text ({token_count} tokens) exceeds {self.MAX_TOKENS} token limit.\n\n"
            f"{action}"
        )

    def _use_local_pdf_path(self, file_path: str) -> bool:
        if self.operation.mode != OperationMode.LOCAL:
            return False
        if _is_unc_path(file_path):
            return False
        return os.path.isfile(file_path)

    def _is_text_read_for_edit(self, file_path: str) -> bool:
        return (
            not self._is_image(file_path)
            and not self._is_pdf(file_path)
            and not file_path.lower().endswith(".ipynb")
            and not self._is_office_doc(file_path)
        )

    def _is_plain_text_candidate(self, file_path: str) -> bool:
        return (
            self._is_text_read_for_edit(file_path)
            and not self._is_binary(file_path)
        )

    async def _read_raw_text_for_edit_state(self, file_path: str) -> Optional[_RawTextState]:
        """Read raw text content for EditFileTool stale-write checks."""
        try:
            res = await self.operation.fs().read_file(file_path)
            if res.code != StatusCode.SUCCESS.code:
                return None
            content = res.data.content if res.data else ""
            if isinstance(content, bytes):
                content = content.decode("utf-8", errors="replace")
            normalized = str(content).replace("\r\n", "\n")
            return _RawTextState(content=normalized, line_count=self._raw_line_count(normalized))
        except Exception:
            return None

    async def _record_read_state(
            self,
            file_path: str,
            mtime_ns: int,
            size_bytes: int,
            is_partial: bool,
            rendered_line_count: int,
    ) -> None:
        # 沙箱模式下本地 os.stat 会失败，mtime_ns 为 0；仍需注册，否则
        # EditFileTool 的预读校验会因为找不到 read_state 而拒绝编辑。
        if not self._is_text_read_for_edit(file_path):
            return

        raw_state = None if is_partial else await self._read_raw_text_for_edit_state(file_path)
        raw_line_count = raw_state.line_count if raw_state else rendered_line_count
        effective_partial = is_partial or raw_line_count > rendered_line_count

        _FILE_READ_REGISTRY[file_path] = _FileReadState(
            mtime_ns=mtime_ns,
            size_bytes=size_bytes,
            is_partial=effective_partial,
            content=raw_state.content if not effective_partial and raw_state else None,
        )

    # ------------------------------------------------------------------
    # Per-type readers
    # ------------------------------------------------------------------

    async def _read_text(self, file_path: str, offset: int, limit: int, apply_size_cap: bool = True) -> str:
        """Read `limit` lines starting after skipping `offset` lines (0-based skip).

        apply_size_cap: apply MAX_SIZE_BYTES check only when reading without an explicit
        user-supplied limit (mirrors TS: maxSizeBytes is passed only when limit is None).
        """
        if apply_size_cap and not _is_unc_path(file_path):
            try:
                byte_len = os.stat(file_path).st_size
                if byte_len > self.MAX_SIZE_BYTES:
                    raise RuntimeError(
                        f"File content ({byte_len // 1024} KB) exceeds maximum allowed size "
                        f"({self.MAX_SIZE_BYTES // 1024} KB). "
                        "Use offset and limit parameters to read specific portions of the file."
                    )
            except OSError:
                pass
        start = max(0, offset) + 1    # convert 0-based skip to 1-indexed start line
        end = start + max(0, limit) - 1
        res = await self.operation.fs().read_file(file_path, line_range=(start, end))
        if res.code != StatusCode.SUCCESS.code:
            raise RuntimeError(res.message)
        content = res.data.content or ""
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")

        if apply_size_cap:
            byte_len = len(content.encode("utf-8", errors="replace"))
            if byte_len > self.MAX_SIZE_BYTES:
                raise RuntimeError(
                    f"File content ({byte_len // 1024} KB) exceeds maximum allowed size "
                    f"({self.MAX_SIZE_BYTES // 1024} KB). "
                    "Use offset and limit parameters to read specific portions of the file."
                )

        tokens = self._estimate_tokens(content)
        if tokens > self.MAX_TOKENS:
            raise MaxFileReadTokenExceededError(tokens, self.MAX_TOKENS)

        rendered = self._cat_n(content)

        # Empty file or offset beyond end-of-file: return an explicit warning string.
        if not content.strip():
            if offset == 0:
                return "Warning: the file exists but the contents are empty."
            else:
                return (
                    f"Warning: the file exists but is shorter than "
                    f"the provided offset ({offset}). "
                    f"The file has {len(content.splitlines())} lines."
                )

        return rendered

    # ------------------------------------------------------------------
    # Office document readers (.docx / .xlsx / .pptx)
    # ------------------------------------------------------------------

    # Map Office document extensions to (import_module_name, pypi_package_name).
    _OFFICE_DOC_PACKAGES: Dict[str, Tuple[str, str]] = {
        ".docx": ("docx", "python-docx"),
        ".xlsx": ("openpyxl", "openpyxl"),
        ".pptx": ("pptx", "python-pptx"),
    }

    async def _read_office_doc(self, file_path: str) -> str:
        """Dispatch Office documents to the appropriate parser and return cat-n text."""
        ext = os.path.splitext(file_path.lower())[1]

        # Pre-parse size check to avoid loading oversized ZIP containers into memory.
        # Office docs use a larger limit than text files because they are ZIP containers
        # with embedded media (images); the real protection is the post-parse token check.
        if not _is_unc_path(file_path):
            try:
                byte_len = os.stat(file_path).st_size
                if byte_len > self.MAX_OFFICE_DOC_SIZE_BYTES:
                    raise RuntimeError(
                        f"Office document '{os.path.basename(file_path)}' ({byte_len // 1024 // 1024} MB) "
                        f"exceeds the maximum allowed size ({self.MAX_OFFICE_DOC_SIZE_BYTES // 1024 // 1024} MB). "
                        "Office documents are parsed in their entirety and cannot be read in portions. "
                        "Consider splitting the document or converting to a smaller format."
                    )
            except OSError:
                pass

        try:
            if ext == ".docx":
                content = await asyncio.to_thread(self._read_docx, file_path)
            elif ext == ".xlsx":
                content = await asyncio.to_thread(self._read_xlsx, file_path)
            elif ext == ".pptx":
                content = await asyncio.to_thread(self._read_pptx, file_path)
            elif ext in (".doc", ".xls", ".ppt"):
                raise RuntimeError(
                    f"Legacy Office format '{ext}' is not supported. "
                    f"Please convert to the modern format ({ext}x) and try again."
                )
            else:
                raise RuntimeError(f"Unsupported Office document format: {ext}")
        except ImportError as exc:
            _info = self._OFFICE_DOC_PACKAGES.get(ext, (exc.name, exc.name))
            pkg_name = _info[1]
            raise RuntimeError(
                f"Reading '{ext}' files requires the '{pkg_name}' package. "
                f"Install with: pip install {pkg_name}"
            ) from exc

        tokens = self._estimate_tokens(content)
        if tokens > self.MAX_TOKENS:
            raise MaxFileReadTokenExceededError(tokens, self.MAX_TOKENS)

        if not content.strip():
            return "Warning: the document exists but the contents are empty."

        return self._cat_n(content)

    @staticmethod
    def _read_docx(file_path: str) -> str:
        """Extract text and tables from a .docx file as Markdown."""
        from docx import Document

        doc = Document(file_path)
        parts: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name if para.style else ""
            if style_name == "Title":
                parts.append(f"# {text}")
            elif style_name.startswith("Heading"):
                try:
                    level = int(style_name.split()[-1])
                except (ValueError, IndexError):
                    level = 1
                parts.append(f"{'#' * (level + 1)} {text}")
            else:
                parts.append(text)

        for table in doc.tables:
            if not table.rows:
                continue
            table_lines = []
            header_cells = [cell.text.replace("|", "\\|").strip() for cell in table.rows[0].cells]
            table_lines.append("| " + " | ".join(header_cells) + " |")
            table_lines.append("| " + " | ".join("---" for _ in header_cells) + " |")
            for row in table.rows[1:]:
                cells = [cell.text.replace("|", "\\|").strip() for cell in row.cells]
                table_lines.append("| " + " | ".join(cells) + " |")
            parts.append("\n".join(table_lines))

        return "\n\n".join(parts)

    @staticmethod
    def _read_xlsx(file_path: str) -> str:
        """Extract data from an .xlsx file as Markdown tables."""
        from openpyxl import load_workbook

        parts: list[str] = []

        wb = load_workbook(file_path, read_only=True, data_only=True)
        try:
            for ws in wb.worksheets:
                table_lines = [f"## {ws.title}"]
                first_row = True
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if not any(cells):
                        continue
                    escaped = [c.replace("|", "\\|") for c in cells]
                    table_lines.append("| " + " | ".join(escaped) + " |")
                    if first_row:
                        table_lines.append("| " + " | ".join("---" for _ in cells) + " |")
                        first_row = False
                parts.append("\n".join(table_lines))
        finally:
            wb.close()

        return "\n\n".join(parts)

    @staticmethod
    def _read_pptx(file_path: str) -> str:
        """Extract text from a .pptx file as Markdown."""
        from pptx import Presentation

        prs = Presentation(file_path)
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {i}\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    table = shape.table
                    if table.rows:
                        header = [cell.text.replace("|", "\\|").strip() for cell in table.rows[0].cells]
                        parts.append("| " + " | ".join(header) + " |")
                        parts.append("| " + " | ".join("---" for _ in header) + " |")
                        for row in table.rows[1:]:
                            cells = [cell.text.replace("|", "\\|").strip() for cell in row.cells]
                            parts.append("| " + " | ".join(cells) + " |")
            parts.append("")

        return "\n".join(parts)

    async def _read_notebook(self, file_path: str) -> str:
        res = await self.operation.fs().read_file(file_path)
        if res.code != StatusCode.SUCCESS.code:
            raise RuntimeError(res.message)
        raw_text = res.data.content or ""
        if isinstance(raw_text, bytes):
            raw_text = raw_text.decode("utf-8", errors="replace")

        byte_len = len(raw_text.encode("utf-8", errors="replace"))
        if byte_len > self.MAX_SIZE_BYTES:
            raise RuntimeError(
                f"Notebook content ({byte_len // 1024} KB) exceeds maximum allowed size "
                f"({self.MAX_SIZE_BYTES // 1024} KB). "
                "Use Bash with jq to inspect specific cells:\n"
                f"  cat \"{file_path}\" | jq '.cells[:20]'        # First 20 cells\n"
                f"  cat \"{file_path}\" | jq '.cells | length'    # Count total cells"
            )

        tokens = self._estimate_tokens(raw_text)
        if tokens > self.MAX_TOKENS:
            raise MaxFileReadTokenExceededError(tokens, self.MAX_TOKENS)

        notebook = json.loads(raw_text)
        cells = notebook.get("cells", [])
        blocks = []
        for idx, cell in enumerate(cells, 1):
            cell_type = cell.get("cell_type", "unknown")
            source = "".join(cell.get("source", []))
            blocks.append(f"## Cell {idx} [{cell_type}]")
            if source:
                blocks.append(source.rstrip("\n"))
            outputs = cell.get("outputs") or []
            if outputs:
                blocks.append("### Outputs")
            for out in outputs:
                text = ""
                if "text" in out:
                    text = "".join(out.get("text", []))
                elif "data" in out and isinstance(out["data"], dict):
                    txt = out["data"].get("text/plain")
                    if isinstance(txt, list):
                        text = "".join(txt)
                    elif isinstance(txt, str):
                        text = txt
                elif "ename" in out and "evalue" in out:
                    text = f"{out.get('ename')}: {out.get('evalue')}"
                if text:
                    blocks.append(text.rstrip("\n"))
        return "\n".join(blocks).strip()

    def _extract_pdf_pages_sync(
            self,
            file_path: str,
            *,
            pages: Optional[str],
            pdf_bytes: Optional[bytes] = None,
    ) -> str:
        try:
            import pdfplumber
        except ImportError as exc:
            raise RuntimeError(
                "Reading PDF files requires the optional dependency 'pdfplumber'."
            ) from exc

        pdf_source = io.BytesIO(pdf_bytes) if pdf_bytes is not None else file_path
        with pdfplumber.open(pdf_source) as pdf:
            total_pages = len(pdf.pages)

            # Require pages param when document is too long to inline safely.
            if not pages and total_pages > self.PDF_AT_MENTION_INLINE_THRESHOLD:
                raise PDFReadError(self._build_pdf_too_many_pages_error(file_path, total_pages))

            result, reason = self._parse_pdf_page_range_with_reason(pages, total_pages)
            if result is None:
                if reason == "out_of_bounds":
                    raise PDFReadError(
                        self._build_pdf_page_range_out_of_bounds_error(file_path, pages, total_pages)
                    )
                raise PDFReadError(self._build_pdf_invalid_page_range_error(str(pages or "")))
            start, end = result

            page_count = end - start + 1
            if page_count > self.PDF_MAX_PAGES_PER_READ:
                raise PDFReadError(self._build_pdf_page_range_too_wide_error(file_path, str(pages), page_count))

            parts = []
            for page_no in range(start, end + 1):
                page_text = pdf.pages[page_no - 1].extract_text() or ""
                parts.append(f"## Page {page_no}\n{page_text}".rstrip())
            return "\n\n".join(parts).strip()

    async def _read_pdf(self, file_path: str, pages: Optional[str]) -> str:
        size_bytes: int = 0
        try:
            size_bytes = os.stat(file_path).st_size
        except OSError:
            size_bytes = 0

        if self.operation.mode == OperationMode.SANDBOX:
            if not size_bytes or size_bytes > self.MAX_PDF_SIZE_BYTES_WITHOUT_PAGES:
                raise PDFReadError(self._build_pdf_sandbox_large_file_error(file_path, pages))

        if not pages and size_bytes > self.MAX_PDF_SIZE_BYTES_WITHOUT_PAGES:
            raise PDFReadError(self._build_pdf_too_large_no_pages_error(file_path, size_bytes))

        if pages and size_bytes > self.MAX_PDF_ABSOLUTE_SIZE_BYTES:
            raise PDFReadError(self._build_pdf_absolute_size_error(file_path, size_bytes, pages))

        started = time.perf_counter()
        async with _PDF_READ_SEMAPHORE:
            if self._use_local_pdf_path(file_path):
                rendered = await asyncio.to_thread(
                    self._extract_pdf_pages_sync,
                    file_path,
                    pages=pages,
                    pdf_bytes=None,
                )
            else:
                res = await self.operation.fs().read_file(file_path, mode="bytes")
                if res.code != StatusCode.SUCCESS.code:
                    raise RuntimeError(res.message)
                content = res.data.content or b""
                if not isinstance(content, bytes):
                    raise RuntimeError("PDF content is not bytes")
                rendered = await asyncio.to_thread(
                    self._extract_pdf_pages_sync,
                    file_path,
                    pages=pages,
                    pdf_bytes=content,
                )

        elapsed_ms = int((time.perf_counter() - started) * 1000)
        logger.info("[ReadFileTool] pdf read completed path=%s pages=%s elapsed_ms=%s", file_path, pages, elapsed_ms)

        tokens = self._estimate_tokens(rendered)
        if tokens > self.MAX_TOKENS:
            raise PDFReadError(self._build_pdf_output_token_error(file_path, pages, tokens))

        return rendered

    @staticmethod
    def _compress_image_bytes(raw: bytes, size: Tuple[int, int], quality: int) -> Optional[bytes]:
        try:
            from PIL import Image
            with Image.open(io.BytesIO(raw)) as img:
                out = io.BytesIO()
                img = img.convert("RGB")
                img.thumbnail(size)
                img.save(out, format="JPEG", quality=quality)
                return out.getvalue()
        except Exception:
            return None

    async def _read_image(self, file_path: str, model_name: str) -> Dict[str, Any]:
        if model_name and not self._is_pdf_supported(model_name):
            return {
                "content": "Current model does not support vision image payload.",
                "multimodal": [],
            }

        res = await self.operation.fs().read_file(file_path, mode="bytes")
        if res.code != StatusCode.SUCCESS.code:
            raise RuntimeError(res.message)
        raw = res.data.content or b""
        if not isinstance(raw, bytes):
            raise RuntimeError("Image content is not bytes")
        if not raw:
            raise RuntimeError(f"Image file is empty: {file_path}")

        _, ext = os.path.splitext(file_path.lower())
        image_type = ext.lstrip(".") or "png"
        dimensions: Optional[str] = None

        # Step 1: standard resize (thumbnail to 1536×1536).
        resized = raw
        try:
            from PIL import Image
            with Image.open(io.BytesIO(raw)) as img:
                detected_format = (img.format or "PNG").lower()
                image_type = detected_format  # prefer format detected from buffer
                dimensions = f"{img.width}x{img.height}"
                img.thumbnail((1536, 1536))
                out = io.BytesIO()
                img.save(out, format=detected_format.upper())
                candidate = out.getvalue()
                if candidate and len(candidate) < len(raw):
                    resized = candidate
        except Exception:
            resized = raw

        # Step 2: token budget check — base64 byte count × 0.125 ≈ tokens.
        estimated_tokens = max(1, int(len(base64.b64encode(resized)) * 0.125))
        if estimated_tokens > self.MAX_TOKENS:
            # Aggressive compression from the same buffer.
            compressed = self._compress_image_bytes(raw, size=(800, 800), quality=40)

            if compressed and int(len(base64.b64encode(compressed)) * 0.125) <= self.MAX_TOKENS:
                resized = compressed
                image_type = "jpeg"
            else:
                # Final fallback: 400×400 JPEG q=20 (mirrors TS Sharp fallback).
                fallback = self._compress_image_bytes(raw, size=(400, 400), quality=20)
                if fallback:
                    resized = fallback
                    image_type = "jpeg"

        mime_type = f"image/{image_type}"
        parts = [
            f"Image file read: {file_path}",
            f"format: {image_type}",
            f"size_bytes: {len(raw)}",
            f"transmitted_size_bytes: {len(resized)}",
        ]
        if dimensions:
            parts.append(f"dimensions: {dimensions}")

        if not self.enable_image_multimodal:
            parts.append(
                "Image bytes are not attached because read_file native image multimodal input is disabled."
            )
            parts.append(
                "If a vision tool is configured, call image_ocr or visual_question_answering with this file path."
            )
            return {
                "content": "\n".join(parts),
                "multimodal": [],
            }

        encoded = base64.b64encode(resized).decode("ascii")
        data_url = f"data:{mime_type};base64,{encoded}"
        parts.append("Image bytes are attached as multimodal input and omitted from this tool result.")
        return {
            "content": "\n".join(parts),
            "multimodal": [
                {
                    "type": "image",
                    "source": "read_file",
                    "source_path": file_path,
                    "mime_type": mime_type,
                    "data_url": data_url,
                }
            ],
        }

    # ------------------------------------------------------------------
    # invoke / stream
    # ------------------------------------------------------------------

    async def invoke(self, inputs: Dict[str, Any], **kwargs) -> ToolOutput:
        file_path = inputs.get("file_path")
        if not file_path:
            return ToolOutput(success=False, error="file_path is required")

        try:
            file_path = _resolve_tool_file_path(self.operation, file_path)
        except ValueError as exc:
            return ToolOutput(success=False, error=str(exc))

        # Blocked device paths — would hang or produce infinite output (errorCode 9).
        if self._is_blocked_device(file_path):
            return ToolOutput(
                success=False,
                error=f"Reading device file '{file_path}' is not allowed.",
            )

        # Binary files are not readable as text (errorCode 4); PDF, images, and Office docs are exempt.
        if not self._is_exempt_from_binary_check(file_path):
            return ToolOutput(
                success=False,
                error=f"Binary files cannot be read as text: '{os.path.basename(file_path)}'.",
            )

        pages: Optional[str] = inputs.get("pages")

        # Validate pages format and range size before any I/O (errorCode 7 / 8).
        if pages is not None and self._is_pdf(file_path):
            pages_text = str(pages)
            if not self._validate_pdf_page_range_format(pages_text):
                return ToolOutput(
                    success=False,
                    error=self._build_pdf_invalid_page_range_error(pages_text),
                )
            # Width check: parse against sys.maxsize so open-ended ranges are bounded.
            parsed = self._parse_pdf_page_range(pages_text, sys.maxsize)
            if parsed is None:
                return ToolOutput(
                    success=False,
                    error=self._build_pdf_invalid_page_range_error(pages_text),
                )
            # Open-ended ranges (e.g. "10-") resolve end to sys.maxsize at parse time,
            # which means we cannot bound the size — reject them upfront (errorCode 8).
            start_pg, end_pg = parsed
            if end_pg >= sys.maxsize or (end_pg - start_pg + 1) > self.PDF_MAX_PAGES_PER_READ:
                requested_count = end_pg - start_pg + 1 if end_pg < sys.maxsize else sys.maxsize
                return ToolOutput(
                    success=False,
                    error=self._build_pdf_page_range_too_wide_error(file_path, str(pages), requested_count),
                )

        offset: int = int(inputs.get("offset", 0))
        raw_limit = inputs.get("limit")
        user_supplied_limit: bool = raw_limit is not None
        limit: int = (
            min(int(raw_limit), self.MAX_LINES_TO_READ)
            if user_supplied_limit
            else self.MAX_LINES_TO_READ
        )
        model_name: str = self._resolve_model_name(kwargs)

        key = (file_path, offset, limit, str(pages or ""))

        # mtime-based dedup: avoid re-reading unchanged files.
        mtime_ns: int = 0
        size_bytes: int = 0
        try:
            _st = os.stat(file_path)
            mtime_ns = _st.st_mtime_ns
            size_bytes = _st.st_size
        except OSError:
            pass

        # previous = self._snapshots.get(key)
        # if previous and mtime_ns and previous.mtime_ns == mtime_ns:
        #     await self._record_read_state(
        #         file_path=file_path,
        #         mtime_ns=mtime_ns,
        #         size_bytes=size_bytes,
        #         is_partial=user_supplied_limit or offset > 0,
        #         rendered_line_count=previous.line_count,
        #     )
        #     return ToolOutput(
        #         success=True,
        #         data={
        #             "content": self.FILE_UNCHANGED_STUB,
        #             "file_path": file_path,
        #             "unchanged": True,
        #             "line_count": previous.line_count,
        #         },
        #     )

        try:
            if self._is_pdf(file_path):
                rendered = await self._read_pdf(file_path, pages)
            elif file_path.lower().endswith(".ipynb"):
                rendered = await self._read_notebook(file_path)
            elif self._is_image(file_path):
                rendered = await self._read_image(file_path, model_name)
            elif self._is_office_doc(file_path):
                rendered = await self._read_office_doc(file_path)
            else:
                # MAX_SIZE_BYTES cap only applies when no explicit limit was given
                # (mirrors TS: maxSizeBytes passed only when limit is undefined).
                rendered = await self._read_text(
                    file_path, offset, limit, apply_size_cap=not user_supplied_limit
                )
        except MaxFileReadTokenExceededError as exc:
            return ToolOutput(success=False, error=str(exc))
        except Exception as exc:
            return ToolOutput(success=False, error=str(exc))

        if isinstance(rendered, dict):
            content = str(rendered.get("content", ""))
            result_data = dict(rendered)
        else:
            content = rendered
            result_data = {"content": content}

        line_count = len(content.splitlines()) if content else 0
        if mtime_ns:
            self._snapshots[key] = _ReadSnapshot(mtime_ns=mtime_ns, line_count=line_count)

        # Populate read registry for EditFileTool pre-read validation.
        await self._record_read_state(
            file_path=file_path,
            mtime_ns=mtime_ns,
            size_bytes=size_bytes,
            is_partial=user_supplied_limit or offset > 0,
            rendered_line_count=line_count,
        )

        return ToolOutput(
            success=True,
            data={
                **result_data,
                "content": content,
                "file_path": file_path,
                "unchanged": False,
                "line_count": line_count,
            },
        )

    async def stream(self, inputs: Dict[str, Any], **kwargs) -> AsyncIterator[Any]:
        if False:
            yield None


# Backward-compatible camelCase aliases for legacy callers.
ReadFileTool.isReadOnly = staticmethod(ReadFileTool.is_read_only)
ReadFileTool.isConcurrencySafe = staticmethod(ReadFileTool.is_concurrency_safe)
ReadFileTool.checkPermissions = staticmethod(ReadFileTool.check_permissions)


class WriteFileTool(Tool):
    """Full-file writer with Claude Write-tool style stale-write protection."""

    MAX_FILE_SIZE: int = 1 * 1024 * 1024 * 1024  # 1 GiB

    def __init__(self, operation: SysOperation, language: str = "cn", agent_id: Optional[str] = None,
                 workspace_path: Optional[str] = None):
        super().__init__(
            build_tool_card("write_file", "WriteFileTool", language, agent_id=agent_id))
        self.operation = operation
        self._agent_id = agent_id or "default"
        self._workspace_path = workspace_path

    @staticmethod
    def _detect_encoding(raw: bytes) -> str:
        return "utf-16-le" if raw[:2] == b"\xff\xfe" else "utf-8"

    async def _read_existing_text(self, file_path: str) -> Tuple[str, str]:
        res = await self.operation.fs().read_file(file_path, mode="bytes")
        if res.code != StatusCode.SUCCESS.code:
            raise OSError(res.message)

        raw = res.data.content if res.data else b""
        if isinstance(raw, str):
            raw = raw.encode("utf-8")

        encoding = self._detect_encoding(raw)
        return raw.decode(encoding, errors="replace"), encoding

    async def invoke(self, inputs: Dict[str, Any], **kwargs) -> ToolOutput:
        path: Optional[str] = inputs.get("file_path")
        content = inputs.get("content")

        if not path:
            return ToolOutput(success=False, error="file_path is required")
        if content is None:
            return ToolOutput(success=False, error="content is required")
        if not isinstance(content, str):
            return ToolOutput(success=False, error="content must be a string")

        try:
            path = _resolve_tool_file_path(self.operation, path)
        except ValueError as exc:
            return ToolOutput(success=False, error=str(exc))

        is_unc = _is_unc_path(path)

        encoding = "utf-8"
        old_content: Optional[str] = None
        operation_type = "create"

        if not is_unc:
            try:
                file_exists = os.path.exists(path)
                if file_exists:
                    if os.path.isdir(path):
                        return ToolOutput(success=False, error=f"Target path is a directory: {path}")

                    stat = os.stat(path)
                    if stat.st_size > self.MAX_FILE_SIZE:
                        return ToolOutput(
                            success=False,
                            error=(
                                f"File is too large ({stat.st_size // (1024 ** 3)} GiB). "
                                "Maximum allowed size is 1 GiB."
                            ),
                        )

                    read_state = _FILE_READ_REGISTRY.get(path)
                    if read_state is not None and read_state.is_partial:
                        return ToolOutput(
                            success=False,
                            error=(
                                f"File was only partially read (offset/limit). "
                                f"Call read_file on '{path}' without offset or limit "
                                "to read the entire file before writing to it."
                            ),
                        )
                    if read_state is None:
                        return ToolOutput(
                            success=False,
                            error=(
                                f"File has not been read yet. "
                                f"Call read_file on '{path}' first before writing to it "
                                "(read the entire file; do not use offset or limit)."
                            ),
                        )

                    old_content, encoding = await self._read_existing_text(path)
                    old_content_lf = old_content.replace("\r\n", "\n")

                    if read_state.mtime_ns != stat.st_mtime_ns or read_state.size_bytes != stat.st_size:
                        content_unchanged = read_state.content is not None and old_content_lf == read_state.content
                        if not content_unchanged:
                            _FILE_READ_REGISTRY.pop(path, None)
                            return ToolOutput(
                                success=False,
                                error=(
                                    "File has been modified since read, either by the user or by a linter. "
                                    "Read it again before attempting to write it."
                                ),
                            )

                    operation_type = "update"
            except OSError as exc:
                return ToolOutput(success=False, error=str(exc))

        res = await self.operation.fs().write_file(
            path,
            content,
            prepend_newline=False,
            create_if_not_exist=True,
            encoding=encoding,
        )
        if res.code != StatusCode.SUCCESS.code:
            return ToolOutput(success=False, error=res.message)

        if not is_unc:
            try:
                stat_after = os.stat(path)
                _FILE_READ_REGISTRY[path] = _FileReadState(
                    mtime_ns=stat_after.st_mtime_ns,
                    size_bytes=stat_after.st_size,
                    is_partial=False,
                    content=content.replace("\r\n", "\n"),
                )
            except OSError:
                _FILE_READ_REGISTRY.pop(path, None)

        _session = kwargs.get("session")
        if _session is not None:
            _base_dir = self._workspace_path or str(pathlib.Path(get_cwd()).expanduser().resolve())
            _history_path = os.path.join(
                _base_dir, ".agent_history", f"file_ops_{self._agent_id}_{_session.get_session_id()}.json"
            )
            await _append_op_history(_history_path, path, "write", old_content, content)

        return ToolOutput(
            success=True,
            data={
                "file_path": path,
                "bytes_written": len(content.encode(encoding, errors="replace")),
                "type": operation_type,
                "created": operation_type == "create",
                "original_file": old_content,
            }
        )

    async def stream(self, inputs: Dict[str, Any], **kwargs) -> AsyncIterator[Any]:
        pass


class EditFileTool(Tool):
    """Enhanced file edit tool implementing the full Edit tool specification.

    Key behaviours:
    - Pre-read validation: file must be read by ReadFileTool before editing
    - External modification detection: mtime_ns + file-size dual check
    - Uniqueness validation: old_string must match exactly once unless replace_all=True
    - Quote tolerance: automatic straight/curly quote fallback when exact match fails
    - XML desanitization: reverses HTML entity encoding via html.unescape()
    - New file creation: old_string='' + non-existent path creates new file
    - EOL preservation: detects and preserves CRLF vs LF on write
    - Trailing whitespace stripping on new_string lines (skipped for .md / .mdx)
    - File size guard: rejects files > 1 GiB
    - Rejects .ipynb files (use NotebookEdit instead)
    - Rejects old_string == new_string (no-op)
    - UNC path handling: skips local filesystem checks for \\\\server or //server paths
    """

    MAX_FILE_SIZE: int = 1 * 1024 * 1024 * 1024  # 1 GiB
    _MD_EXTENSIONS: frozenset = frozenset({".md", ".mdx"})
    _STRAIGHT_TO_CURLY: dict = str.maketrans({'"': '\u201c', '\u201c': '"', '\u201d': '"',
                                               "'": '\u2018', '\u2018': "'", '\u2019': "'"})

    def __init__(self, operation: SysOperation, language: str = "cn", agent_id: Optional[str] = None,
        workspace_path: Optional[str] = None):
        super().__init__(build_tool_card("edit_file", "EditFileTool", language, agent_id=agent_id))
        self.operation = operation
        self._agent_id = agent_id or "default"
        self._workspace_path = workspace_path
    # ------------------------------------------------------------------
    # Static helpers
    # ------------------------------------------------------------------

    # Mirrors TS DESANITIZATIONS table in FileEditTool/utils.ts.
    # Claude cannot output certain XML tokens (they're sanitized by the API), so it
    # emits compressed abbreviations that the Edit tool must reverse.
    _DESANITIZATIONS: dict = {
        "<fnr>": "<function_results>",
        "<n>": "<name>",
        "</n>": "</name>",
        "<o>": "<output>",
        "</o>": "</output>",
        "<e>": "<error>",
        "</e>": "</error>",
        "<s>": "<system>",
        "</s>": "</system>",
        "<r>": "<result>",
        "</r>": "</result>",
        "< META_START >": "<META_START>",
        "< META_END >": "<META_END>",
        "< EOT >": "<EOT>",
        "< META >": "<META>",
        "< SOS >": "<SOS>",
        "\n\nH:": "\n\nHuman:",
        "\n\nA:": "\n\nAssistant:",
    }

    @staticmethod
    def _detect_eol(raw: bytes) -> str:
        return "\r\n" if b"\r\n" in raw else "\n"

    @staticmethod
    def _detect_encoding(raw: bytes) -> str:
        return "utf-16-le" if raw[:2] == b"\xff\xfe" else "utf-8"

    @classmethod
    def _desanitize(cls, value: str) -> str:
        result = html.unescape(value)
        for source, target in cls._DESANITIZATIONS.items():
            result = result.replace(source, target)
        return result

    @staticmethod
    def _strip_trailing_whitespace(value: str) -> str:
        parts = value.splitlines(keepends=True)
        if not parts:
            return value.rstrip()

        stripped_parts: List[str] = []
        for part in parts:
            body = part
            line_ending = ""
            if part.endswith("\r\n"):
                body = part[:-2]
                line_ending = "\r\n"
            elif part.endswith(("\n", "\r")):
                body = part[:-1]
                line_ending = part[-1]
            stripped_parts.append(body.rstrip() + line_ending)
        return "".join(stripped_parts)

    @staticmethod
    def _find_similar_paths(file_path: str, max_results: int = 5) -> List[str]:
        directory = os.path.dirname(file_path) or "."
        target_base = os.path.splitext(os.path.basename(file_path))[0].lower()
        try:
            names = os.listdir(directory)
        except OSError:
            return []

        similar: List[str] = []
        for name in names:
            candidate_base = os.path.splitext(name)[0].lower()
            if candidate_base == target_base or target_base in candidate_base or candidate_base in target_base:
                similar.append(os.path.join(directory, name))
            if len(similar) >= max_results:
                break
        return similar

    @staticmethod
    def _normalize_quotes(value: str) -> str:
        return (
            value.replace("\u2018", "'")
            .replace("\u2019", "'")
            .replace("\u201c", '"')
            .replace("\u201d", '"')
        )

    @staticmethod
    def _is_opening_quote_context(chars: List[str], index: int) -> bool:
        if index == 0:
            return True
        return chars[index - 1] in {" ", "\t", "\n", "\r", "(", "[", "{", "\u2014", "\u2013"}

    @classmethod
    def _apply_curly_double_quotes(cls, value: str) -> str:
        chars = list(value)
        result: List[str] = []
        for index, char in enumerate(chars):
            if char == '"':
                result.append("\u201c" if cls._is_opening_quote_context(chars, index) else "\u201d")
            else:
                result.append(char)
        return "".join(result)

    @classmethod
    def _apply_curly_single_quotes(cls, value: str) -> str:
        chars = list(value)
        result: List[str] = []
        for index, char in enumerate(chars):
            if char != "'":
                result.append(char)
                continue
            prev_char = chars[index - 1] if index > 0 else ""
            next_char = chars[index + 1] if index < len(chars) - 1 else ""
            if prev_char.isalpha() and next_char.isalpha():
                result.append("\u2019")
            else:
                result.append("\u2018" if cls._is_opening_quote_context(chars, index) else "\u2019")
        return "".join(result)

    @classmethod
    def _preserve_quote_style(cls, old_str: str, actual_old_str: str, new_str: str) -> str:
        if old_str == actual_old_str:
            return new_str

        result = new_str
        if "\u201c" in actual_old_str or "\u201d" in actual_old_str:
            result = cls._apply_curly_double_quotes(result)
        if "\u2018" in actual_old_str or "\u2019" in actual_old_str:
            result = cls._apply_curly_single_quotes(result)
        return result

    def _try_quote_variants(self, content: str, old_str: str) -> Optional[str]:
        """Return a quote-substituted variant of old_str that matches content, or None."""
        normalized_content = self._normalize_quotes(content)
        normalized_old = self._normalize_quotes(old_str)
        index = normalized_content.find(normalized_old)
        if index == -1:
            return None
        return content[index:index + len(old_str)]

    async def _read_existing_text(self, file_path: str) -> Tuple[str, bytes]:
        res = await self.operation.fs().read_file(file_path, mode="bytes")
        if res.code != StatusCode.SUCCESS.code:
            raise OSError(res.message)

        raw = res.data.content if res.data else b""
        if isinstance(raw, str):
            raw = raw.encode("utf-8")

        encoding = self._detect_encoding(raw)
        return raw.decode(encoding, errors="replace"), raw

    # ------------------------------------------------------------------
    # invoke
    # ------------------------------------------------------------------

    async def invoke(self, inputs: Dict[str, Any], **kwargs) -> ToolOutput:  # noqa: C901
        file_path: Optional[str] = inputs.get("file_path")
        old_str: str = inputs.get("old_string", "")
        new_str: Optional[str] = inputs.get("new_string")
        replace_all: bool = bool(inputs.get("replace_all", False))

        if not file_path:
            return ToolOutput(success=False, error="file_path is required")
        if new_str is None:
            return ToolOutput(success=False, error="new_string is required")

        try:
            file_path = _resolve_tool_file_path(self.operation, file_path)
        except ValueError as exc:
            return ToolOutput(success=False, error=str(exc))


        # Reject Jupyter notebooks — use NotebookEdit instead.
        if file_path.lower().endswith(".ipynb"):
            return ToolOutput(
                success=False,
                error="Cannot edit .ipynb files with this tool. Use NotebookEdit instead.",
            )

        # Reject no-op edits.
        if old_str == new_str:
            return ToolOutput(
                success=False,
                error="old_string and new_string are identical; no changes would be made.",
            )

        is_unc = _is_unc_path(file_path)
        if is_unc:
            file_exists = True
        elif self.operation.mode == OperationMode.SANDBOX:
            # 沙箱模式下 os.path.exists 只能看到主容器的文件系统，
            # 实际文件在 jiuwenbox sidecar 中。用沙箱 fs 探测真实存在性。
            try:
                _probe = await self.operation.fs().read_file(file_path, mode="bytes")
                file_exists = _probe.code == StatusCode.SUCCESS.code
            except Exception:
                file_exists = False
        else:
            file_exists = os.path.exists(file_path)

        # ---- New file creation (old_string == '') --------------------------------
        if old_str == "":
            if file_exists and not is_unc:
                try:
                    existing_content, _ = await self._read_existing_text(file_path)
                except OSError as exc:
                    return ToolOutput(success=False, error=str(exc))
                if existing_content.strip() != "":
                    return ToolOutput(
                        success=False,
                        error="Cannot create new file - file already exists.",
                    )
            write_res = await self.operation.fs().write_file(
                file_path, new_str, prepend_newline=False, create_if_not_exist=True
            )
            if write_res.code != StatusCode.SUCCESS.code:
                return ToolOutput(success=False, error=f"Create failed: {write_res.message}")
            # Register the new file as read so subsequent edits don't require a re-read.
            try:
                _st = os.stat(file_path)
                _FILE_READ_REGISTRY[file_path] = _FileReadState(
                    mtime_ns=_st.st_mtime_ns,
                    size_bytes=_st.st_size,
                    is_partial=False,
                    content=new_str.replace("\r\n", "\n"),
                )
            except OSError:
                pass
            return ToolOutput(success=True, data={"file_path": file_path, "replacements": 0, "created": True})

        # ---- File must exist for non-empty old_string ----------------------------
        if not file_exists:
            similar = self._find_similar_paths(file_path)
            hint = f" Similar paths: {similar}" if similar else ""
            return ToolOutput(
                success=False,
                error=f"File not found: '{file_path}'.{hint}",
            )

        # ---- File size guard -----------------------------------------------------
        if not is_unc:
            _stat = None
            try:
                _stat = os.stat(file_path)
            except OSError as exc:
                # 沙箱模式下文件在 jiuwenbox sidecar 中，本地 os.stat 会失败。
                # 优雅降级到 (0, 0)，外部修改检测已有 content 比对回退路径。
                if self.operation.mode != OperationMode.SANDBOX:
                    return ToolOutput(success=False, error=str(exc))

            if _stat is not None and _stat.st_size > self.MAX_FILE_SIZE:
                return ToolOutput(
                    success=False,
                    error=(
                        f"File is too large ({_stat.st_size // (1024 ** 3)} GiB). "
                        "Maximum allowed size is 1 GiB."
                    ),
                )

            current_mtime = _stat.st_mtime_ns if _stat else 0
            current_size = _stat.st_size if _stat else 0
        else:
            current_mtime = 0
            current_size = 0

        # ---- Pre-read validation -------------------------------------------------
        read_state = _FILE_READ_REGISTRY.get(file_path)
        if read_state is not None and read_state.is_partial:
            return ToolOutput(
                success=False,
                error=(
                    f"File was only partially read (offset/limit). "
                    f"Call read_file on '{file_path}' without offset or limit "
                    "to read the entire file before editing."
                ),
            )
        if read_state is None:
            return ToolOutput(
                success=False,
                error=(
                    f"File must be read before editing. "
                    f"Call read_file on '{file_path}' first "
                    "(read the entire file; do not use offset or limit)."
                ),
            )

        # ---- External modification check (timestamp + size dual check) -----------
        if not is_unc and (read_state.mtime_ns != current_mtime or read_state.size_bytes != current_size):
            content_unchanged = False
            if read_state.content is not None:
                try:
                    compare_content, _ = await self._read_existing_text(file_path)
                    compare_content = compare_content.replace("\r\n", "\n")
                    content_unchanged = compare_content == read_state.content
                except OSError:
                    content_unchanged = False

            if not content_unchanged:
                _FILE_READ_REGISTRY.pop(file_path, None)
                return ToolOutput(
                    success=False,
                    error=(
                        f"'{file_path}' has been modified externally since it was last read. "
                        "Re-read the file before editing."
                    ),
                )

        # ---- Read raw bytes to detect EOL style ----------------------------------
        try:
            content, raw = await self._read_existing_text(file_path)
        except OSError as exc:
            return ToolOutput(success=False, error=str(exc))

        eol = self._detect_eol(raw)

        # Normalise to LF internally for matching and replacement.
        content_lf = content.replace("\r\n", "\n")

        # ---- XML desanitization + trailing whitespace stripping ------------------
        old_str_clean = self._desanitize(old_str).replace("\r\n", "\n")
        _, ext = os.path.splitext(file_path.lower())
        preserve_md = ext in self._MD_EXTENSIONS
        new_str_clean = self._desanitize(new_str).replace("\r\n", "\n")
        if not preserve_md:
            new_str_clean = self._strip_trailing_whitespace(new_str_clean)

        # ---- Matching + quote tolerance ------------------------------------------
        match_str = old_str_clean
        if match_str not in content_lf:
            variant = self._try_quote_variants(content_lf, old_str_clean)
            if variant is None:
                return ToolOutput(
                    success=False,
                    error=f"old_string not found in '{file_path}'.",
                )
            match_str = variant
            new_str_clean = self._preserve_quote_style(old_str_clean, match_str, new_str_clean)

        # ---- Uniqueness validation -----------------------------------------------
        count = content_lf.count(match_str)
        if not replace_all and count > 1:
            return ToolOutput(
                success=False,
                error=(
                    f"old_string matches {count} times in '{file_path}'. "
                    "Provide more surrounding context to make it unique, "
                    "or set replace_all=true to replace every occurrence."
                ),
            )

        # ---- Perform replacement -------------------------------------------------
        if replace_all:
            new_content_lf = content_lf.replace(match_str, new_str_clean)
            replaced = count
        else:
            new_content_lf = content_lf.replace(match_str, new_str_clean, 1)
            replaced = 1

        # Restore original EOL style.
        new_content = new_content_lf.replace("\n", eol) if eol == "\r\n" else new_content_lf

        # ---- Write back ----------------------------------------------------------
        write_res = await self.operation.fs().write_file(
            file_path, new_content, prepend_newline=False
        )
        if write_res.code != StatusCode.SUCCESS.code:
            return ToolOutput(success=False, error=f"Write failed: {write_res.message}")

        # Update registry so the next edit doesn't require a re-read.
        try:
            _st2 = os.stat(file_path)
            _FILE_READ_REGISTRY[file_path] = _FileReadState(
                mtime_ns=_st2.st_mtime_ns,
                size_bytes=_st2.st_size,
                is_partial=False,
                content=new_content_lf,
            )
        except OSError:
            _FILE_READ_REGISTRY.pop(file_path, None)

        _session = kwargs.get("session")
        if _session is not None:
            _base_dir = self._workspace_path or str(pathlib.Path(get_cwd()).expanduser().resolve())
            _history_path = os.path.join(
                _base_dir, ".agent_history", f"file_ops_{self._agent_id}_{_session.get_session_id()}.json"
            )
            await _append_op_history(_history_path, file_path, "edit", content, new_content)

        return ToolOutput(
            success=True,
            data={"file_path": file_path, "replacements": replaced},
        )

    async def stream(self, inputs: Dict[str, Any], **kwargs) -> AsyncIterator[Any]:
        if False:
            yield None


class GlobTool(Tool):
    DEFAULT_MAX_RESULTS: int = 100

    def __init__(self, operation: SysOperation, language: str = "cn", agent_id: Optional[str] = None):
        super().__init__(
            build_tool_card("glob", "GlobTool", language, agent_id=agent_id))
        self.operation = operation

    def _resolve_search_path(self, path: Optional[str]) -> str:
        if path:
            return _resolve_tool_file_path(self.operation, path)
        return str(pathlib.Path(get_cwd()).expanduser().resolve())

    def _relativize_paths(self, paths: List[str], base_path: str) -> List[str]:
        relative_paths: List[str] = []
        for item in paths:
            try:
                relative_paths.append(os.path.relpath(item, base_path))
            except ValueError:
                relative_paths.append(item)
        return relative_paths

    async def invoke(self, inputs: Dict[str, Any], **kwargs) -> ToolOutput:
        pattern = inputs.get("pattern")
        if not pattern:
            return ToolOutput(success=False, error="pattern is required")

        try:
            path = self._resolve_search_path(inputs.get("path"))
        except ValueError as exc:
            return ToolOutput(success=False, error=str(exc))

        started_at = time.perf_counter()

        res = await self.operation.fs().search_files(path, pattern)
        if res.code != StatusCode.SUCCESS.code:
            return ToolOutput(success=False, error=res.message)

        matching_files = [item.path for item in res.data.matching_files] if res.data else []
        truncated = len(matching_files) > self.DEFAULT_MAX_RESULTS
        limited_files = matching_files[:self.DEFAULT_MAX_RESULTS]
        filenames = self._relativize_paths(limited_files, path)
        duration_ms = int((time.perf_counter() - started_at) * 1000)

        return ToolOutput(
            success=True,
            data={
                "durationMs": duration_ms,
                "numFiles": len(filenames),
                "filenames": filenames,
                "truncated": truncated,
                "matching_files": limited_files,
                "count": len(filenames),
            }
        )

    async def stream(self, inputs: Dict[str, Any], **kwargs) -> AsyncIterator[Any]:
        pass


class ListDirTool(Tool):

    def __init__(self, operation: SysOperation, language: str = "cn", agent_id: Optional[str] = None):
        super().__init__(
            build_tool_card("list_files", "ListDirTool", language, agent_id=agent_id))
        self.operation = operation

    async def invoke(self, inputs: Dict[str, Any], **kwargs) -> ToolOutput:
        path = inputs.get("path", ".")
        show_hidden = inputs.get("show_hidden", False)

        files_res = await self.operation.fs().list_files(path)
        dirs_res = await self.operation.fs().list_directories(path)

        if files_res.code != StatusCode.SUCCESS.code:
            return ToolOutput(success=False, error=f"列出文件失败: {files_res.message}")
        if dirs_res.code != StatusCode.SUCCESS.code:
            return ToolOutput(success=False, error=f"列出目录失败: {dirs_res.message}")

        files = [item.name for item in files_res.data.list_items] if files_res.data else []
        dirs = [item.name for item in dirs_res.data.list_items] if dirs_res.data else []

        if not show_hidden:
            files = [f for f in files if not f.startswith('.')]
            dirs = [d for d in dirs if not d.startswith('.')]

        files.sort()
        dirs.sort()

        return ToolOutput(
            success=True,
            data={
                "files": files,
                "dirs": dirs
            }
        )

    async def stream(self, inputs: Dict[str, Any], **kwargs) -> AsyncIterator[Any]:
        pass


class GrepTool(Tool):
    DEFAULT_HEAD_LIMIT: int = 250
    MAX_COLUMNS: int = 500
    VCS_DIRECTORIES_TO_EXCLUDE: Tuple[str, ...] = (".git", ".svn", ".hg", ".bzr", ".jj", ".sl")

    def __init__(self, operation: SysOperation, language: str = "cn", agent_id: Optional[str] = None):
        super().__init__(
            build_tool_card("grep", "GrepTool", language, agent_id=agent_id))
        self.operation = operation

    @staticmethod
    def _shell_quote(value: Any) -> str:
        text = str(value)
        if os.name == "nt":
            return "'" + text.replace("'", "''") + "'"
        return shlex.quote(text)

    @staticmethod
    def _as_bool(value: Any, default: bool = False) -> bool:
        if value is None:
            return default
        if isinstance(value, bool):
            return value
        if isinstance(value, (int, float)):
            return bool(value)
        text = str(value).strip().lower()
        if text in {"1", "true", "yes", "on"}:
            return True
        if text in {"0", "false", "no", "off"}:
            return False
        return default

    @staticmethod
    def _as_int(value: Any, default: Optional[int] = None) -> Optional[int]:
        if value is None or value == "":
            return default
        try:
            return int(value)
        except (TypeError, ValueError):
            return default

    @classmethod
    def _apply_head_limit(
            cls, items: List[str], limit: Optional[int], offset: int = 0
    ) -> Tuple[List[str], Optional[int]]:
        offset = max(0, offset)
        if limit == 0:
            return items[offset:], None

        effective_limit = limit if limit is not None else cls.DEFAULT_HEAD_LIMIT
        sliced = items[offset: offset + effective_limit]
        was_truncated = len(items) - offset > effective_limit
        return sliced, effective_limit if was_truncated else None

    @staticmethod
    def _split_glob_patterns(glob_value: Optional[str]) -> List[str]:
        if not glob_value:
            return []

        glob_patterns: List[str] = []
        for raw_pattern in str(glob_value).split():
            if "{" in raw_pattern and "}" in raw_pattern:
                glob_patterns.append(raw_pattern)
            else:
                glob_patterns.extend(part for part in raw_pattern.split(",") if part)
        return glob_patterns

    def _resolve_search_path(self, path: Optional[str]) -> str:
        if path:
            return _resolve_tool_file_path(self.operation, path)
        return str(pathlib.Path(get_cwd()).expanduser().resolve())

    def _build_rg_command(
            self,
            *,
            pattern: str,
            path: str,
            glob: Optional[str],
            output_mode: str,
            context_before: Optional[int],
            context_after: Optional[int],
            context_c: Optional[int],
            context: Optional[int],
            show_line_numbers: bool,
            case_insensitive: bool,
            file_type: Optional[str],
            multiline: bool,
    ) -> str:
        parts: List[str] = [
            "rg",
            "--hidden",
            "--color=never",
            "--max-columns",
            str(self.MAX_COLUMNS),
        ]

        if multiline:
            parts.extend(["-U", "--multiline-dotall"])
        if case_insensitive:
            parts.append("-i")

        if output_mode == "files_with_matches":
            parts.append("-l")
        elif output_mode == "count":
            parts.append("-c")
        elif show_line_numbers:
            parts.append("-n")

        if output_mode == "content":
            if context is not None:
                parts.extend(["-C", str(context)])
            elif context_c is not None:
                parts.extend(["-C", str(context_c)])
            else:
                if context_before is not None:
                    parts.extend(["-B", str(context_before)])
                if context_after is not None:
                    parts.extend(["-A", str(context_after)])

        for directory in self.VCS_DIRECTORIES_TO_EXCLUDE:
            parts.extend(["--glob", self._shell_quote(f"!{directory}")])

        if file_type:
            parts.extend(["--type", self._shell_quote(file_type)])

        for glob_pattern in self._split_glob_patterns(glob):
            parts.extend(["--glob", self._shell_quote(glob_pattern)])

        if pattern.startswith("-"):
            parts.extend(["-e", self._shell_quote(pattern)])
        else:
            parts.append(self._shell_quote(pattern))

        parts.append(self._shell_quote(path))
        return " ".join(parts)

    def _build_select_string_command(
            self,
            *,
            pattern: str,
            path: str,
            glob: Optional[str],
            output_mode: str,
            context_before: Optional[int],
            context_after: Optional[int],
            context_c: Optional[int],
            context: Optional[int],
            case_insensitive: bool,
    ) -> str:
        """Build a PowerShell Select-String command as fallback when rg is unavailable on Windows."""
        sq = self._shell_quote

        # Expand {ts,tsx} brace syntax — PS -like doesn't support it
        glob_patterns = self._split_glob_patterns(glob)
        expanded_globs: List[str] = []
        for p in glob_patterns:
            m = re.match(r'^(.*)\{([^}]+)\}(.*)$', p)
            if m:
                expanded_globs.extend(f"{m.group(1)}{alt}{m.group(3)}" for alt in m.group(2).split(","))
            else:
                expanded_globs.append(p)

        # Context lines: -C / context takes priority over -B / -A
        effective_c = context if context is not None else context_c
        ctx_b = effective_c if effective_c is not None else (context_before or 0)
        ctx_a = effective_c if effective_c is not None else (context_after or 0)

        # VCS exclusion regex for -notmatch: (\\|/)(\.git|\.svn|...)(\\|/|$)
        vcs_alts = "|".join(d.replace(".", r"\.") for d in self.VCS_DIRECTORIES_TO_EXCLUDE)
        vcs_pat = sq(r"(\\|/)(" + vcs_alts + r")(\\|/|$)")

        # --- Pipeline stages ---
        pipeline: List[str] = []

        # Stage 1: file enumeration + VCS pruning
        is_file = os.path.isfile(path)
        if is_file:
            pipeline.append(f"Get-Item -LiteralPath {sq(path)}")
        else:
            pipeline.append(f"Get-ChildItem -LiteralPath {sq(path)} -Recurse -File")
            pipeline.append(f"Where-Object {{ $_.FullName -notmatch {vcs_pat} }}")

        # Stage 2: glob filter (skip for single-file input)
        if expanded_globs and not is_file:
            conds = " -or ".join(f"$_.Name -like {sq(p)}" for p in expanded_globs)
            pipeline.append(f"Where-Object {{ {conds} }}")

        # Stage 3: Select-String
        # Default is case-insensitive; -CaseSensitive enables exact matching
        cs_flag = " -CaseSensitive" if not case_insensitive else ""
        ctx_flag = f" -Context {ctx_b},{ctx_a}" if output_mode == "content" and (ctx_b or ctx_a) else ""
        pipeline.append(f"Select-String -Pattern {sq(pattern)}{cs_flag}{ctx_flag}")

        # Stage 4: output formatting to match rg/grep line format
        if output_mode == "files_with_matches":
            pipeline.append("Select-Object -ExpandProperty Path -Unique")
        elif output_mode == "count":
            pipeline.append("Group-Object Path | ForEach-Object { \"$($_.Name):$($_.Count)\" }")
        elif ctx_b or ctx_a:
            # Manually expand PreContext / PostContext arrays into filepath:linenum:content lines
            pipeline.append(
                "ForEach-Object {"
                " $m=$_; $p=$m.Context.PreContext.Length;"
                " for($i=0;$i-lt$p;$i++){ \"$($m.Path):$([int]$m.LineNumber-$p+$i):$($m.Context.PreContext[$i])\" };"
                " \"$($m.Path):$($m.LineNumber):$($m.Line)\";"
                " for($i=0;$i-lt$m.Context.PostContext.Length;$i++)"
                "{ \"$($m.Path):$([int]$m.LineNumber+1+$i):$($m.Context.PostContext[$i])\" }"
                " }"
            )
        else:
            pipeline.append("ForEach-Object { \"$($_.Path):$($_.LineNumber):$($_.Line)\" }")

        return "$ErrorActionPreference='SilentlyContinue'; " + " | ".join(pipeline)

    def _build_grep_command(
            self,
            *,
            pattern: str,
            path: str,
            glob: Optional[str],
            output_mode: str,
            context_before: Optional[int],
            context_after: Optional[int],
            context_c: Optional[int],
            context: Optional[int],
            show_line_numbers: bool,
            case_insensitive: bool,
            multiline: bool,
    ) -> Optional[str]:
        if multiline:
            return None

        parts: List[str] = ["grep", "-R", "--binary-files=without-match"]

        for directory in self.VCS_DIRECTORIES_TO_EXCLUDE:
            parts.append(f"--exclude-dir={self._shell_quote(directory)}")

        if case_insensitive:
            parts.append("-i")

        if output_mode == "files_with_matches":
            parts.append("-l")
        elif output_mode == "count":
            parts.append("-c")
        elif show_line_numbers:
            parts.append("-n")

        if output_mode == "content":
            if context is not None:
                parts.extend(["-C", str(context)])
            elif context_c is not None:
                parts.extend(["-C", str(context_c)])
            else:
                if context_before is not None:
                    parts.extend(["-B", str(context_before)])
                if context_after is not None:
                    parts.extend(["-A", str(context_after)])

        for glob_pattern in self._split_glob_patterns(glob):
            parts.append(f"--include={self._shell_quote(glob_pattern)}")

        parts.extend([self._shell_quote(pattern), self._shell_quote(path)])
        return " ".join(parts)

    @staticmethod
    def _extract_file_path_from_line(line: str, mode: str) -> Optional[str]:
        if not line:
            return None
        if mode == "files_with_matches":
            return line
        if mode == "count":
            if ":" not in line:
                return None
            return line.rsplit(":", 1)[0]

        match = re.match(r"^(.*?):(\d+|[-]+):(.*)$", line)
        if match:
            return match.group(1)
        return line.split(":", 1)[0] if ":" in line else None

    @staticmethod
    def _relativize_line(line: str, base_path: str, mode: str) -> str:
        file_path = GrepTool._extract_file_path_from_line(line, mode)
        if not file_path:
            return line

        try:
            relative_path = os.path.relpath(file_path, base_path)
        except ValueError:
            return line

        if mode == "files_with_matches":
            return relative_path

        prefix = file_path + ":"
        if line.startswith(prefix):
            return relative_path + ":" + line[len(prefix):]
        return line

    def _build_structured_output(
            self,
            *,
            stdout: str,
            stderr: str,
            exit_code: int,
            output_mode: str,
            head_limit: Optional[int],
            offset: int,
            base_path: str,
    ) -> Dict[str, Any]:
        raw_lines = [line for line in stdout.splitlines() if line.strip()]
        limited_lines, applied_limit = self._apply_head_limit(raw_lines, head_limit, offset)
        final_lines = [self._relativize_line(line, base_path, output_mode) for line in limited_lines]
        content = "\n".join(final_lines)

        data: Dict[str, Any] = {
            "stdout": content,
            "stderr": stderr,
            "exit_code": exit_code,
            "mode": output_mode,
            "appliedOffset": offset if offset > 0 else None,
            "appliedLimit": applied_limit,
        }

        if output_mode == "content":
            data.update({
                "content": content,
                "filenames": [],
                "numFiles": 0,
                "numLines": len(final_lines),
                "count": len(final_lines),
            })
            return data

        if output_mode == "count":
            total_matches = 0
            file_count = 0
            for line in final_lines:
                if ":" not in line:
                    continue
                count_str = line.rsplit(":", 1)[1]
                try:
                    total_matches += int(count_str)
                    file_count += 1
                except ValueError:
                    continue
            data.update({
                "content": content,
                "filenames": [],
                "numFiles": file_count,
                "numMatches": total_matches,
                "count": total_matches,
            })
            return data

        data.update({
            "filenames": final_lines,
            "numFiles": len(final_lines),
            "count": len(final_lines),
        })
        return data

    async def invoke(self, inputs: Dict[str, Any], **kwargs) -> ToolOutput:
        pattern = inputs.get("pattern")
        if not pattern:
            return ToolOutput(success=False, error="pattern is required")

        try:
            path = self._resolve_search_path(inputs.get("path"))
        except ValueError as exc:
            return ToolOutput(success=False, error=str(exc))

        output_mode = str(inputs.get("output_mode") or "content")
        if output_mode not in {"content", "files_with_matches", "count"}:
            return ToolOutput(success=False, error="output_mode must be one of: content, files_with_matches, count")

        ignore_case = self._as_bool(inputs.get("-i", inputs.get("ignore_case", False)))
        show_line_numbers = self._as_bool(inputs.get("-n", True), default=True)
        context_before = self._as_int(inputs.get("-B"))
        context_after = self._as_int(inputs.get("-A"))
        context_c = self._as_int(inputs.get("-C"))
        context = self._as_int(inputs.get("context"))
        head_limit = self._as_int(inputs.get("head_limit"))
        offset = self._as_int(inputs.get("offset"), 0) or 0
        multiline = self._as_bool(inputs.get("multiline", False))
        glob = inputs.get("glob")
        file_type = inputs.get("type")

        has_context_controls = any(
            value is not None for value in [context_before, context_after, context_c, context]
        )
        if output_mode != "content" and has_context_controls:
            context_before = None
            context_after = None
            context_c = None
            context = None

        if shutil.which("rg"):
            cmd = self._build_rg_command(
                pattern=str(pattern),
                path=path,
                glob=glob,
                output_mode=output_mode,
                context_before=context_before,
                context_after=context_after,
                context_c=context_c,
                context=context,
                show_line_numbers=show_line_numbers,
                case_insensitive=ignore_case,
                file_type=file_type,
                multiline=multiline,
            )
        elif os.name == "nt":
            if file_type:
                return ToolOutput(success=False, error="type filter requires ripgrep (rg) to be installed")
            if multiline:
                return ToolOutput(success=False, error="multiline search requires ripgrep (rg) to be installed")
            cmd = self._build_select_string_command(
                pattern=str(pattern),
                path=path,
                glob=glob,
                output_mode=output_mode,
                context_before=context_before,
                context_after=context_after,
                context_c=context_c,
                context=context,
                case_insensitive=ignore_case,
            )
        else:
            if file_type:
                return ToolOutput(success=False, error="type filter requires ripgrep (rg) to be installed")
            cmd = self._build_grep_command(
                pattern=str(pattern),
                path=path,
                glob=glob,
                output_mode=output_mode,
                context_before=context_before,
                context_after=context_after,
                context_c=context_c,
                context=context,
                show_line_numbers=show_line_numbers,
                case_insensitive=ignore_case,
                multiline=multiline,
            )
            if cmd is None:
                return ToolOutput(success=False, error="multiline search requires ripgrep (rg) to be installed")

        shell_type = "powershell" if os.name == "nt" else "auto"
        res = await self.operation.shell().execute_cmd(cmd, timeout=30, shell_type=shell_type)

        if res.code != StatusCode.SUCCESS.code:
            return ToolOutput(success=False, error=res.message)

        stdout = res.data.stdout if res.data else ""
        stderr = res.data.stderr if res.data else ""
        exit_code = res.data.exit_code if res.data else -1
        success = (exit_code in [0, 1])

        return ToolOutput(
            success=success,
            data=self._build_structured_output(
                stdout=stdout,
                stderr=stderr,
                exit_code=exit_code,
                output_mode=output_mode,
                head_limit=head_limit,
                offset=offset,
                base_path=path if os.path.isdir(path) else os.path.dirname(path) or ".",
            ),
            error=stderr if not success else None
        )

    async def stream(self, inputs: Dict[str, Any], **kwargs) -> AsyncIterator[Any]:
        pass
